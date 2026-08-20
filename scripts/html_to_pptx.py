"""Convert presentation-titres-tresor.html content into a PowerPoint deck."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
# Hors du dépôt git (sibling de ekonzo-web)
OUT_DIR = ROOT.parent / "presentations"
OUT = OUT_DIR / "presentation-titres-tresor.pptx"
LOGO = ROOT / "public" / "logo.webp"

# Charte Gouvernement RDC
BLEU = RGBColor(0x00, 0x95, 0xC9)
JAUNE = RGBColor(0xFF, 0xF2, 0x4B)
ROUGE = RGBColor(0xDB, 0x38, 0x32)
BLEU_INST = RGBColor(0x17, 0x48, 0x99)
BLEU_FONCE = RGBColor(0x17, 0x41, 0x8A)
NOIR = RGBColor(0x32, 0x32, 0x30)
MUTED = RGBColor(0x5A, 0x5A, 0x58)
BLANC = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = RGBColor(0xF5, 0xF8, 0xFB)
CARD_Y = RGBColor(0xFF, 0xFC, 0xE8)
CARD_R = RGBColor(0xFD, 0xF0, 0xEF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run(run, text, size=14, bold=False, color=NOIR, font="Arial"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_textbox(slide, left, top, width, height, text, size=14, bold=False, color=NOIR, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size, bold, color)
    return box


def no_shadow(shape):
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def rect(slide, left, top, width, height, fill: RGBColor):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    no_shadow(shape)
    return shape


def stripe_footer(slide, label: str, cover=False):
    y = SLIDE_H - Inches(0.55)
    w = Inches(0.18)
    band = Inches(0.1)
    x0 = Inches(0.45)
    for i, c in enumerate((BLEU, JAUNE, ROUGE)):
        rect(slide, x0, y + Inches(0.02) + band * i, w, band, c)
    color = RGBColor(0xAA, 0xCC, 0xDD) if cover else MUTED
    add_textbox(slide, Inches(0.85), y, Inches(11.5), Inches(0.35), label, size=10, color=color)


def brand_bar(slide, meta: str, cover=False):
    bg = BLEU_FONCE if cover else BLANC
    if cover:
        rect(slide, 0, 0, SLIDE_W, SLIDE_H, BLEU_FONCE)
        rect(slide, 0, 0, SLIDE_W, Inches(1.05), BLEU_INST)
    else:
        rect(slide, 0, 0, SLIDE_W, Inches(1.05), BLANC)
        # bottom border line
        rect(slide, 0, Inches(1.05), SLIDE_W, Inches(0.02), RGBColor(0xD6, 0xE0, 0xEE))

    # logo
    if LOGO.exists():
        try:
            from PIL import Image

            img = Image.open(LOGO).convert("RGBA")
            buf = BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            pic = slide.shapes.add_picture(buf, Inches(0.45), Inches(0.22), height=Inches(0.65))
            no_shadow(pic)
        except Exception:
            pass

    meta_color = RGBColor(0xCC, 0xDD, 0xEE) if cover else MUTED
    add_textbox(
        slide,
        Inches(9.2),
        Inches(0.28),
        Inches(3.7),
        Inches(0.6),
        meta,
        size=11,
        color=meta_color,
        align=PP_ALIGN.RIGHT,
    )


def eyebrow(slide, text: str, cover=False):
    color = BLEU if cover else BLEU_INST
    add_textbox(slide, Inches(0.55), Inches(1.25), Inches(12), Inches(0.35), text.upper(), size=11, bold=True, color=color)


def title(slide, text: str, cover=False, y=Inches(1.55)):
    color = BLANC if cover else BLEU_FONCE
    add_textbox(slide, Inches(0.55), y, Inches(12.2), Inches(0.9), text, size=28 if not cover else 32, bold=True, color=color)


def lead(slide, text: str, cover=False, y=Inches(2.35)):
    color = RGBColor(0xE0, 0xEE, 0xF8) if cover else NOIR
    add_textbox(slide, Inches(0.55), y, Inches(12), Inches(0.9), text, size=15, color=color)


def card(slide, left, top, width, height, heading, body_lines, accent="n"):
    fills = {"n": CARD_BG, "y": CARD_Y, "r": CARD_R}
    accents = {"n": BLEU_INST, "y": JAUNE, "r": ROUGE}
    rect(slide, left, top, width, height, fills.get(accent, CARD_BG))
    # left accent bar
    rect(slide, left, top, Inches(0.08), height, accents.get(accent, BLEU_INST))
    add_textbox(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.3), Inches(0.35), heading, size=14, bold=True, color=BLEU_FONCE)
    box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.45), width - Inches(0.3), height - Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in body_lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(4)
        run = p.add_run()
        set_run(run, line, 12, False, NOIR)


def add_table(slide, left, top, width, rows, col_widths=None):
    rows_n = len(rows)
    cols_n = len(rows[0])
    table_shape = slide.shapes.add_table(rows_n, cols_n, left, top, width, Inches(0.42 * rows_n))
    no_shadow(table_shape)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            is_header = r == 0
            set_run(run, val, 11 if not is_header else 12, is_header, BLANC if is_header else NOIR)
            fill = BLEU_FONCE if is_header else (RGBColor(0xF0, 0xF4, 0xF8) if r % 2 else BLANC)
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def build():
    prs = new_prs()

    # 1 Cover
    s = blank(prs)
    brand_bar(s, "République Démocratique\ndu Congo", cover=True)
    eyebrow(s, "Présentation de projet", cover=True)
    title(s, "ekonzo — souscription digitale aux titres du Trésor", cover=True)
    lead(
        s,
        "Théorie des Bons & Obligations du Trésor, processus actuel d’achat, puis digitalisation via ekonzo — avec ses contraintes réglementaires.",
        cover=True,
    )
    stripe_footer(s, "Ministère des Finances · Plateforme ekonzo", cover=True)

    # 2 Agenda
    s = blank(prs)
    brand_bar(s, "Sommaire")
    eyebrow(s, "Plan de la présentation")
    title(s, "Au programme")
    items = [
        "01  Théorie : titres, BT, OT",
        "02  Acheteurs actuels & profils cibles",
        "03  Processus actuel d’achat (BT & OT)",
        "04  Limites, digitalisation ekonzo",
        "05  Contraintes qui demeurent",
        "06  Parcours, paiements & bénéfices",
    ]
    box = s.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(11), Inches(3.8))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        run = p.add_run()
        set_run(run, it, 18, True, BLEU_FONCE)
    stripe_footer(s, "ekonzo · Ministère des Finances")

    # 3 Titre du Trésor
    s = blank(prs)
    brand_bar(s, "Partie 1 · Théorie")
    eyebrow(s, "Théorie")
    title(s, "Qu’est-ce qu’un titre du Trésor ?")
    lead(s, "C’est un prêt que vous faites à l’État. Le Ministère des Finances s’engage à rembourser à une date précise, avec un rendement défini à l’émission.")
    cards = [
        ("Émetteur", ["L’État congolais, via le Ministère des Finances (Trésor public)."], "n"),
        ("Garantie", ["Engagement souverain. Titres dématérialisés, inscrits en comptes-titres."], "y"),
        ("Devises", ["Émissions possibles en CDF ou en USD."], "r"),
        ("Cadre légal", ["Décret n°18/025 du 11 juin 2018 — modalités d’émission et de remboursement."], "n"),
    ]
    for i, (h, b, a) in enumerate(cards):
        card(s, Inches(0.45 + i * 3.15), Inches(3.4), Inches(3.0), Inches(2.4), h, b, a)
    stripe_footer(s, "Titres publics · RDC")

    # 4 BT
    s = blank(prs)
    brand_bar(s, "Partie 1 · Bons du Trésor")
    eyebrow(s, "Théorie")
    title(s, "Les Bons du Trésor (BT)")
    lead(s, "Titres à court terme pour financer les besoins de trésorerie de l’État. Idéals pour une épargne de quelques mois — cœur de l’offre ekonzo aujourd’hui.")
    card(
        s,
        Inches(0.5),
        Inches(3.35),
        Inches(5.9),
        Inches(2.8),
        "Caractéristiques",
        [
            "• Maturité ≤ 12 mois",
            "• Titres négociables et dématérialisés",
            "• Rendement par taux d’escompte",
            "• Intérêts et capital à l’échéance",
        ],
        "n",
    )
    card(
        s,
        Inches(6.7),
        Inches(3.35),
        Inches(5.9),
        Inches(2.8),
        "Pour l’investisseur",
        [
            "• Placement sûr, horizon court",
            "• Montants accessibles (dès 10 000 CDF)",
            "• Souscription via canaux agréés",
            "• Garantie de l’État congolais",
        ],
        "y",
    )
    stripe_footer(s, "Théorie · Bons du Trésor")

    # 5 OT
    s = blank(prs)
    brand_bar(s, "Partie 1 · Obligations du Trésor")
    eyebrow(s, "Théorie")
    title(s, "Les Obligations du Trésor (OT)")
    lead(s, "Titres à moyen / long terme (> 12 mois). L’investisseur perçoit des coupons (intérêts périodiques) jusqu’au remboursement du capital.")
    card(
        s,
        Inches(0.5),
        Inches(3.35),
        Inches(5.9),
        Inches(2.8),
        "Caractéristiques",
        [
            "• Maturité supérieure à 12 mois",
            "• Versement de coupons selon le calendrier",
            "• Capital remboursé à l’échéance",
            "• Outil de financement structurel de l’État",
        ],
        "n",
    )
    card(
        s,
        Inches(6.7),
        Inches(3.35),
        Inches(5.9),
        Inches(2.8),
        "Pour l’investisseur",
        [
            "• Revenus périodiques prévisibles",
            "• Horizon d’épargne plus long",
            "• Profil institutionnel ou particulier avisé",
            "• Même garantie souveraine que les BT",
        ],
        "r",
    )
    stripe_footer(s, "Théorie · Obligations du Trésor")

    # 6 Compare BT OT
    s = blank(prs)
    brand_bar(s, "Partie 1 · Synthèse")
    eyebrow(s, "Théorie")
    title(s, "BT vs OT — synthèse")
    add_table(
        s,
        Inches(0.55),
        Inches(2.5),
        Inches(12.2),
        [
            ["Critère", "Bon du Trésor (BT)", "Obligation du Trésor (OT)"],
            ["Durée", "Jusqu’à 12 mois", "Plus de 12 mois"],
            ["Rendement", "Taux d’escompte (à l’échéance)", "Coupons d’intérêts périodiques"],
            ["Objectif État", "Trésorerie", "Financement moyen / long terme"],
            ["Profil", "Court terme, liquidité", "Horizon long, revenus"],
            ["Sécurité", "Garantie de l’État congolais", "Garantie de l’État congolais"],
        ],
        [Inches(2.4), Inches(4.9), Inches(4.9)],
    )
    stripe_footer(s, "Fin de la synthèse BT / OT")

    # 7 Qui achète
    s = blank(prs)
    brand_bar(s, "Partie 1 · Marché actuel")
    eyebrow(s, "Aujourd’hui · Avant digitalisation")
    title(s, "Qui achète vraiment auprès du Trésor ?")
    lead(
        s,
        "Aujourd’hui, la grande majorité des BT et OT est souscrite par les banques (et des assurances / institutionnels via le circuit bancaire) — pas par le grand public.",
        y=Inches(2.3),
    )
    card(
        s,
        Inches(0.45),
        Inches(3.25),
        Inches(6.0),
        Inches(3.3),
        "Les banques",
        [
            "Acheteurs directs aux adjudications · compte propre",
            "Besoin : placer la liquidité excédentaire.",
            "Motivation : rendement sûr, actifs liquides.",
            "Exemple : surplus de trésorerie → BT 91 jours.",
        ],
        "n",
    )
    card(
        s,
        Inches(6.7),
        Inches(3.25),
        Inches(6.0),
        Inches(3.3),
        "Les assurances (et assimilés)",
        [
            "Institutionnels · souvent via banque / dépositaire",
            "Besoin : actifs « État » pour engagements longs.",
            "Motivation : sécurité, coupons (OT).",
            "Exemple : réserves techniques sur OT 3–5 ans.",
        ],
        "y",
    )
    stripe_footer(s, "Marché primaire actuel · banques & institutionnels")

    # 8 Exemples banque vs assurance
    s = blank(prs)
    brand_bar(s, "Partie 1 · Exemples institutionnels")
    eyebrow(s, "Aujourd’hui · Cas types")
    title(s, "Exemples : banque vs assurance")
    card(
        s,
        Inches(0.4),
        Inches(2.5),
        Inches(4.0),
        Inches(3.5),
        "Banque — compte propre (BT)",
        ["Trésorerie bancaire", "Excédent de dépôts à placer.", "Cas : surplus → BT 28 ou 91 jours."],
        "n",
    )
    card(
        s,
        Inches(4.6),
        Inches(2.5),
        Inches(4.0),
        Inches(3.5),
        "Banque — pour ses clients",
        ["Intermédiation", "Ordres clients via banque.", "Cas : importateur 2 M USD de BT."],
        "y",
    )
    card(
        s,
        Inches(8.8),
        Inches(2.5),
        Inches(4.0),
        Inches(3.5),
        "Assurance — OT",
        ["Réserves techniques", "Placement long avec coupons.", "Cas : OT 3–5 ans via banque."],
        "r",
    )
    add_textbox(
        s,
        Inches(0.55),
        Inches(6.15),
        Inches(12),
        Inches(0.5),
        "À retenir : le Trésor vend surtout à des acteurs qui ont déjà accès à l’adjudication et au compte-titres. Le particulier reste en marge — c’est le vide qu’ekonzo vise.",
        size=12,
        bold=True,
        color=BLEU_FONCE,
    )
    stripe_footer(s, "Profils qui achètent aujourd’hui auprès du Trésor")

    # 9 Profils BT
    s = blank(prs)
    brand_bar(s, "Partie 1 · Profils cibles BT")
    eyebrow(s, "Cibles ekonzo · Peu présents aujourd’hui")
    title(s, "Profils grand public (BT) — exemples cibles")
    lead(s, "Horizon court : ces profils existent en demande, mais passent rarement par le circuit actuel.", y=Inches(2.35))
    card(s, Inches(0.4), Inches(3.2), Inches(4.0), Inches(3.2), "Marie, 34 ans", ["Enseignante · Kinshasa", "Cas : 800 000 CDF sur BT 6 mois pour les frais scolaires."], "n")
    card(s, Inches(4.6), Inches(3.2), Inches(4.0), Inches(3.2), "Jean, 41 ans", ["Commerçant importateur", "Cas : 5 000 000 CDF pendant 3 mois entre deux commandes."], "y")
    card(s, Inches(8.8), Inches(3.2), Inches(4.0), Inches(3.2), "Grace, 29 ans", ["Diaspora · Bruxelles", "Cas : 1 000 USD sur un BT en dollar via canal digital."], "r")
    stripe_footer(s, "Profils BT à démocratiser")

    # 10 Profils OT
    s = blank(prs)
    brand_bar(s, "Partie 1 · Profils cibles OT")
    eyebrow(s, "Cibles ekonzo · Complément long terme")
    title(s, "Profils grand public / PME (OT) — exemples cibles")
    lead(s, "Horizon moyen / long : aujourd’hui surtout assurances & banques ; demain aussi particuliers avertis et PME.", y=Inches(2.35))
    card(s, Inches(0.4), Inches(3.2), Inches(4.0), Inches(3.2), "Paul, 52 ans", ["Cadre · études des enfants", "Cas : 10 000 USD sur OT 3 ans, coupons semestriels."], "n")
    card(s, Inches(4.6), Inches(3.2), Inches(4.0), Inches(3.2), "SARL « Kivu Trade »", ["PME · réserves d’entreprise", "Cas : OT 2–5 ans plutôt que tout laisser en cash."], "y")
    card(s, Inches(8.8), Inches(3.2), Inches(4.0), Inches(3.2), "Rappel marché actuel", ["Assurance · banque", "Cas : volumes OT encore portés par les institutionnels."], "r")
    stripe_footer(s, "OT · institutionnels aujourd’hui · grand public demain")

    # 11 Processus actuel
    s = blank(prs)
    brand_bar(s, "Partie 2 · Processus actuel")
    eyebrow(s, "Terrain · Aujourd’hui")
    title(s, "Comment on achète un BT ou une OT aujourd’hui")
    lead(s, "Pour les Bons comme pour les Obligations, le circuit reste centré sur la banque : l’adjudication est ouverte aux banques ; le particulier passe presque toujours par un guichet.")
    steps = [
        ("01", "Banque", "Se rendre en agence / être client"),
        ("02", "Compte-titres", "Ouverture obligatoire (loi 22/069)"),
        ("03", "Dossier", "KYC papier, formulaires, signatures"),
        ("04", "Ordre", "Instruction d’achat BT ou OT"),
        ("05", "Adjudication", "La banque place l’ordre au Trésor"),
    ]
    for i, (n, h, d) in enumerate(steps):
        x = Inches(0.4 + i * 2.55)
        rect(s, x, Inches(3.4), Inches(2.4), Inches(2.3), CARD_BG)
        rect(s, x, Inches(3.4), Inches(2.4), Inches(0.08), BLEU)
        add_textbox(s, x + Inches(0.15), Inches(3.55), Inches(2.1), Inches(0.4), n, size=20, bold=True, color=BLEU_INST)
        add_textbox(s, x + Inches(0.15), Inches(4.05), Inches(2.1), Inches(0.4), h, size=14, bold=True, color=BLEU_FONCE)
        add_textbox(s, x + Inches(0.15), Inches(4.5), Inches(2.1), Inches(1.0), d, size=11, color=NOIR)
    add_textbox(
        s,
        Inches(0.55),
        Inches(5.9),
        Inches(12),
        Inches(0.5),
        "Le titre est ensuite inscrit sur le compte-titres tenu par la banque (teneur de compte agréé BCC) — pas sur un « compte ekonzo ».",
        size=12,
        color=MUTED,
    )
    stripe_footer(s, "Circuit bancaire traditionnel · BT & OT")

    # 12 Freins
    s = blank(prs)
    brand_bar(s, "Partie 2 · Limites")
    eyebrow(s, "Terrain · Contres")
    title(s, "Pourquoi ce processus freine l’accès")
    card(
        s,
        Inches(0.45),
        Inches(2.5),
        Inches(6.0),
        Inches(4.0),
        "Pour le citoyen",
        [
            "• Obligation d’avoir un compte-titres bancaire",
            "• Déplacements, files d’attente, horaires d’agence",
            "• Dossier papier long (KYC, convention, ordres)",
            "• Seuils et frais bancaires parfois dissuasifs",
            "• Peu d’information claire sur les émissions",
            "• Difficile pour la diaspora et les non-bancarisés",
        ],
        "r",
    )
    card(
        s,
        Inches(6.7),
        Inches(2.5),
        Inches(6.0),
        Inches(4.0),
        "Pour le Trésor / l’État",
        [
            "• Base d’investisseurs limitée aux clients bancaires",
            "• Faible mobilisation de l’épargne populaire",
            "• Peu de traçabilité digitale côté particulier",
            "• Image d’un produit « réservé aux initiés »",
            "• Canal de collecte peu adapté au Mobile Money",
        ],
        "y",
    )
    stripe_footer(s, "Freins à l’inclusion financière")

    # 13 Digitaliser
    s = blank(prs)
    brand_bar(s, "Partie 3 · Digitalisation")
    eyebrow(s, "Projet ekonzo")
    title(s, "Ce que nous voulons digitaliser")
    lead(s, "ekonzo = canal digital de collecte et d’accompagnement : le citoyen souscrit en ligne ; la banque partenaire reste le teneur de compte-titres.")
    card(
        s,
        Inches(0.45),
        Inches(3.3),
        Inches(6.0),
        Inches(3.2),
        "À digitaliser",
        [
            "• Inscription & authentification (téléphone + OTP)",
            "• KYC à distance (pièce + selfie)",
            "• Signature électronique de la convention",
            "• Catalogue des émissions BT (puis OT)",
            "• Paiement Mobile Money / virement",
            "• Suivi du portefeuille & dossier CIF",
        ],
        "n",
    )
    card(
        s,
        Inches(6.7),
        Inches(3.3),
        Inches(6.0),
        Inches(3.2),
        "Ce qui reste bancaire / Trésor",
        [
            "• Tenue du compte-titres (banque agréée BCC)",
            "• Participation à l’adjudication",
            "• Règlement–livraison des titres",
            "• Calendrier et conditions d’émission (Ministère)",
        ],
        "y",
    )
    stripe_footer(s, "Digitaliser le parcours citoyen, pas remplacer la banque")

    # 14 Avant / après
    s = blank(prs)
    brand_bar(s, "Partie 3 · Avant / après")
    eyebrow(s, "Projet ekonzo")
    title(s, "Processus actuel vs parcours ekonzo")
    add_table(
        s,
        Inches(0.55),
        Inches(2.45),
        Inches(12.2),
        [
            ["Étape", "Aujourd’hui (banque)", "Avec ekonzo"],
            ["Accès", "Agence, clientèle bancaire", "Web / mobile, 24h/24"],
            ["Identité", "KYC papier en agence", "KYC digital (pièce + selfie)"],
            ["Compte-titres", "Ouverture longue en banque", "Convention en ligne → banque partenaire"],
            ["Ordre BT / OT", "Formulaire / guichet", "Choix d’émission + paiement en ligne"],
            ["Paiement", "Compte bancaire / caisse", "MoMo (Airtel, Orange, M-Pesa) ou virement"],
            ["Suivi", "Relevés bancaires", "Tableau de bord investisseur"],
        ],
        [Inches(2.4), Inches(4.9), Inches(4.9)],
    )
    stripe_footer(s, "Même titre d’État · parcours citoyen modernisé")

    # 15 Contraintes
    s = blank(prs)
    brand_bar(s, "Partie 3 · Contraintes")
    eyebrow(s, "Projet · Cadre légal")
    title(s, "Contraintes que ekonzo ne peut pas lever seul")
    lead(s, "La digitalisation accélère le parcours, mais certaines règles restent intangibles.")
    card(
        s,
        Inches(0.45),
        Inches(3.25),
        Inches(6.0),
        Inches(3.3),
        "Contraintes majeures",
        [
            "• Compte-titres obligatoire (loi n°22/069)",
            "• ekonzo n’est pas teneur de compte-titres",
            "• L’adjudication reste Trésor / banques",
            "• Convention écrite + instruction BCC",
        ],
        "r",
    )
    card(
        s,
        Inches(6.7),
        Inches(3.25),
        Inches(6.0),
        Inches(3.3),
        "Autres points de vigilance",
        [
            "• Dépendance à un partenariat bancaire",
            "• KYC / AML obligatoires",
            "• Arrêté ministériel utile (décret 18/025)",
            "• Éducation financière BT vs OT",
            "• Disponibilité réseau / Mobile Money",
        ],
        "n",
    )
    stripe_footer(s, "Transparence sur ce qui change… et ce qui ne change pas")

    # 16 Parcours
    s = blank(prs)
    brand_bar(s, "Partie 3 · Parcours ekonzo")
    eyebrow(s, "Projet ekonzo")
    title(s, "Parcours digital cible")
    lead(s, "De l’inscription à la détention du titre, sans file d’attente en agence.")
    steps = [
        ("01", "Compte", "Téléphone + OTP"),
        ("02", "KYC", "Pièce d’identité & selfie"),
        ("03", "Convention", "Compte-titres (banque partenaire)"),
        ("04", "Souscription", "BT (puis OT) + paiement"),
        ("05", "Titre", "Adjudication via la banque"),
    ]
    for i, (n, h, d) in enumerate(steps):
        x = Inches(0.4 + i * 2.55)
        rect(s, x, Inches(3.35), Inches(2.4), Inches(2.2), CARD_BG)
        rect(s, x, Inches(3.35), Inches(2.4), Inches(0.08), BLEU)
        add_textbox(s, x + Inches(0.15), Inches(3.5), Inches(2.1), Inches(0.4), n, size=20, bold=True, color=BLEU_INST)
        add_textbox(s, x + Inches(0.15), Inches(4.0), Inches(2.1), Inches(0.4), h, size=14, bold=True, color=BLEU_FONCE)
        add_textbox(s, x + Inches(0.15), Inches(4.45), Inches(2.1), Inches(0.9), d, size=11, color=NOIR)
    add_textbox(
        s,
        Inches(0.55),
        Inches(5.8),
        Inches(12),
        Inches(0.5),
        "Le dossier client (CIF) est transmis à la banque pour ouverture / tenue du compte-titres.",
        size=12,
        color=MUTED,
    )
    stripe_footer(s, "Parcours digital · conformité bancaire respectée")

    # 17 Publics & paiements
    s = blank(prs)
    brand_bar(s, "Partie 3 · Cibles & paiements")
    eyebrow(s, "Projet ekonzo")
    title(s, "Pour qui ? Comment payer ?")
    card(s, Inches(0.4), Inches(2.45), Inches(4.0), Inches(2.0), "Particuliers", ["Épargne sur titres d’État via Mobile Money — sans parcours 100 % papier."], "n")
    card(s, Inches(4.6), Inches(2.45), Inches(4.0), Inches(2.0), "Professionnels", ["Placer la trésorerie sur des instruments publics réglementés."], "y")
    card(s, Inches(8.8), Inches(2.45), Inches(4.0), Inches(2.0), "Diaspora", ["Accéder aux BT / OT sans être physiquement en agence à Kinshasa."], "r")
    pays = [
        ("Airtel Money", "Paiement via prompt USSD."),
        ("Orange Money", "Souscription en quelques secondes."),
        ("M-Pesa", "Paiement Mobile Money sécurisé."),
        ("Virement", "Transfert bancaire avec référence."),
    ]
    for i, (h, d) in enumerate(pays):
        card(s, Inches(0.4 + i * 3.2), Inches(4.7), Inches(3.05), Inches(1.6), h, [d], "n")
    stripe_footer(s, "Inclusion financière · Canaux digitaux")

    # 18 Bénéfices
    s = blank(prs)
    brand_bar(s, "Partie 3 · Impact")
    eyebrow(s, "Projet ekonzo")
    title(s, "Bénéfices attendus")
    card(
        s,
        Inches(0.45),
        Inches(2.5),
        Inches(6.0),
        Inches(3.8),
        "Pour le citoyen",
        [
            "• Moins de déplacements, parcours clair",
            "• Épargne sûre, traçable, suivie en ligne",
            "• Paiement Mobile Money familier",
            "• Accès progressif BT puis OT",
        ],
        "n",
    )
    card(
        s,
        Inches(6.7),
        Inches(2.5),
        Inches(6.0),
        Inches(3.8),
        "Pour l’État",
        [
            "• Élargissement de la base d’investisseurs",
            "• Mobilisation de l’épargne nationale",
            "• Canal digital auditable (CIF, KYC)",
            "• Modernisation Trésor–citoyen, dans le cadre légal",
        ],
        "y",
    )
    stripe_footer(s, "Impact citoyen & Trésor public")

    # 19 Conclusion
    s = blank(prs)
    brand_bar(s, "Gouvernement de la République", cover=True)
    eyebrow(s, "Conclusion", cover=True)
    title(s, "Digitaliser l’accès, respecter le cadre bancaire", cover=True)
    lead(
        s,
        "ekonzo modernise le parcours d’achat des BT et OT, tout en conservant le compte-titres et l’adjudication chez la banque partenaire — pour un Trésor plus accessible, sans court-circuiter la loi.",
        cover=True,
    )
    add_textbox(
        s,
        Inches(0.55),
        Inches(4.5),
        Inches(12),
        Inches(0.6),
        "Merci de votre attention · Questions ?",
        size=18,
        bold=True,
        color=BLANC,
    )
    stripe_footer(s, "ekonzo · Ministère des Finances · Kinshasa, RDC", cover=True)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Saved {OUT} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
